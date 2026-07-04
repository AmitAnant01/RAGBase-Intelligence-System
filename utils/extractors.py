"""
Extracts plain text from many file types: PDF (+ scanned/OCR), DOCX, PPTX,
XLSX, CSV, TXT, MD, HTML, JSON.
"""

import os
import json
import csv
import io


def extract_text(file_path: str) -> str:
    ext = file_path.rsplit(".", 1)[-1].lower()

    if ext == "pdf":
        return _extract_pdf(file_path)
    elif ext == "docx":
        return _extract_docx(file_path)
    elif ext == "pptx":
        return _extract_pptx(file_path)
    elif ext in ("xlsx", "xlsm"):
        return _extract_xlsx(file_path)
    elif ext == "csv":
        return _extract_csv(file_path)
    elif ext in ("txt", "md"):
        return _extract_txt(file_path)
    elif ext in ("html", "htm"):
        return _extract_html(file_path)
    elif ext == "json":
        return _extract_json(file_path)
    elif ext in ("png", "jpg", "jpeg", "webp", "bmp"):
        return _extract_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    text_parts = []
    needs_ocr_pages = []

    for idx, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""
        if len(page_text.strip()) < 20:
            needs_ocr_pages.append(idx)
        text_parts.append(page_text)

    # Fallback to OCR for pages that came back (near) empty - i.e. scanned PDFs
    if needs_ocr_pages:
        ocr_text = _ocr_pdf_pages(file_path, needs_ocr_pages)
        for idx, txt in ocr_text.items():
            if txt.strip():
                text_parts[idx] = txt

    return "\n".join(text_parts)


def _ocr_pdf_pages(file_path: str, page_indices) -> dict:
    """Best-effort OCR for scanned pages. Silently no-ops if OCR deps are missing."""
    results = {}
    try:
        import pytesseract
        from PIL import Image
        from pypdf import PdfReader

        # Try to rasterize via pdf2image; if unavailable, skip OCR gracefully.
        try:
            from pdf2image import convert_from_path
        except ImportError:
            return results

        images = convert_from_path(file_path)
        for idx in page_indices:
            if idx < len(images):
                results[idx] = pytesseract.image_to_string(images[idx])
    except Exception:
        # OCR engine (tesseract binary) or poppler not installed on this machine.
        # We degrade gracefully - the raw (possibly empty) text is kept instead.
        pass
    return results


def _extract_image(file_path: str) -> str:
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(file_path)
        return pytesseract.image_to_string(img)
    except Exception:
        return ""


def _extract_docx(file_path: str) -> str:
    import docx

    doc = docx.Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)

    return "\n".join(paragraphs)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text)
    return "\n".join(parts)


def _extract_xlsx(file_path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_csv(file_path: str) -> str:
    parts = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            if row:
                parts.append(" | ".join(row))
    return "\n".join(parts)


def _extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_html(file_path: str) -> str:
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _extract_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)
